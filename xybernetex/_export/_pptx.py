"""
PPTX renderer — converts an Artifact to a PowerPoint presentation.
Requires: pip install xybernetex[pptx]   (python-pptx)

Layout:
  Slide 1  — Title slide (artifact title + type)
  Slide N  — One slide per H2 section; bullet points from lists and
              paragraphs; tables rendered as PPTX tables.
"""
from __future__ import annotations

from pathlib import Path
from typing import TYPE_CHECKING

if TYPE_CHECKING:
    from xybernetex._models import Artifact

_MISSING = (
    "python-pptx is required for PPTX export.\n"
    "Install it with:  pip install xybernetex[pptx]"
)

# Brand palette
_DARK   = "1A1A2E"
_ACCENT = "E94560"
_LIGHT  = "F5F5F5"
_WHITE  = "FFFFFF"
_GREY   = "666666"


try:
    from pptx import Presentation as _Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    _Presentation = None  # type: ignore


def render_pptx(artifact: "Artifact", path: str) -> str:
    """
    Write *artifact* to a PowerPoint presentation at *path*.

    If *path* is a directory the file is named after the artifact title.
    Returns the path of the written file.
    """
    if _Presentation is None:
        raise ImportError(_MISSING)

    from xybernetex._export._markdown import parse, inline_text
    from xybernetex._models import _safe_filename

    dest = _resolve_path(path, artifact.title, ".pptx")
    dest.parent.mkdir(parents=True, exist_ok=True)

    prs = _Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # completely blank

    # ── Slide 1: Title ─────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _fill_background(slide, _DARK)

    _add_text_box(
        slide,
        artifact.title,
        left=Inches(1), top=Inches(2.5), width=Inches(11.33), height=Inches(1.5),
        font_size=Pt(40), bold=True, color=_WHITE, align="left",
    )
    _add_text_box(
        slide,
        artifact.artifact_type.replace("_", " ").title(),
        left=Inches(1), top=Inches(4.1), width=Inches(8), height=Inches(0.5),
        font_size=Pt(16), bold=False, color=_ACCENT, align="left",
    )
    # Accent bar
    _add_rect(slide, left=Inches(1), top=Inches(2.2), width=Inches(1.5), height=Inches(0.06), color=_ACCENT)

    # ── Content slides ─────────────────────────────────────────────────────────
    blocks = parse(artifact.content)

    # Group blocks by H1/H2 boundary → one slide per section
    sections: list[dict] = []
    current: dict = {"title": artifact.title, "blocks": []}

    for block in blocks:
        if block["type"] == "heading" and block["level"] <= 2:
            if current["blocks"]:
                sections.append(current)
            current = {"title": block["text"], "blocks": []}
        else:
            current["blocks"].append(block)
    if current["blocks"] or (not sections and blocks):
        sections.append(current)

    for section in sections:
        _add_content_slide(prs, blank_layout, section["title"], section["blocks"])

    prs.save(str(dest))
    return str(dest)


# ── Slide builder ──────────────────────────────────────────────────────────────

def _add_content_slide(prs, layout, title: str, blocks: list[dict]) -> None:
    slide = prs.slides.add_slide(layout)
    _fill_background(slide, _LIGHT)

    # Title bar background strip
    _add_rect(slide, left=0, top=0, width=prs.slide_width, height=Inches(1.15), color=_DARK)

    # Accent side bar
    _add_rect(slide, left=0, top=0, width=Inches(0.12), height=prs.slide_height, color=_ACCENT)

    # Title text
    _add_text_box(
        slide, title,
        left=Inches(0.3), top=Inches(0.2), width=Inches(12.5), height=Inches(0.75),
        font_size=Pt(24), bold=True, color=_WHITE, align="left",
    )

    # Body area
    body_top  = Inches(1.3)
    body_left = Inches(0.5)
    body_w    = Inches(12.3)
    body_h    = Inches(5.8)

    # Check if section has a table — render it specially
    table_blocks = [b for b in blocks if b["type"] == "table"]
    text_blocks  = [b for b in blocks if b["type"] != "table"]

    # Text column
    if text_blocks:
        lines = _blocks_to_bullet_lines(text_blocks)
        if lines:
            txBox = slide.shapes.add_textbox(body_left, body_top, body_w if not table_blocks else Inches(5.5), body_h)
            tf = txBox.text_frame
            tf.word_wrap = True
            _populate_text_frame(tf, lines)

    # Table (first one only per slide to keep it readable)
    if table_blocks:
        tbl_block = table_blocks[0]
        _add_pptx_table(
            slide, tbl_block,
            left=Inches(6.2) if text_blocks and _blocks_to_bullet_lines(text_blocks) else body_left,
            top=body_top,
            width=Inches(6.5) if text_blocks else body_w,
            height=body_h,
        )


def _blocks_to_bullet_lines(blocks: list[dict]) -> list[tuple[str, int]]:
    """Convert blocks to (text, indent_level) tuples for the text frame."""
    from xybernetex._export._markdown import inline_text
    lines: list[tuple[str, int]] = []
    for b in blocks:
        if b["type"] == "paragraph":
            lines.append((inline_text(b["text"]), 0))
        elif b["type"] in ("bullet_list", "ordered_list"):
            for item in b["items"]:
                lines.append((inline_text(item), 1))
        elif b["type"] == "heading":
            lines.append((b["text"].upper(), 0))
        elif b["type"] == "code":
            for l in b["text"].splitlines()[:6]:  # cap code preview
                lines.append((l, 1))
    return lines


def _populate_text_frame(tf, lines: list[tuple[str, int]]) -> None:
    first = True
    for text, level in lines:
        if first:
            para = tf.paragraphs[0]
            first = False
        else:
            para = tf.add_paragraph()
        para.level = level
        run = para.add_run()
        run.text = ("• " if level > 0 else "") + text
        run.font.size = Pt(13 if level == 0 else 11)
        run.font.color.rgb = RGBColor.from_string(_DARK if level == 0 else _GREY)
        run.font.bold = (level == 0)


def _add_pptx_table(slide, block: dict, left, top, width, height) -> None:
    headers  = block["headers"]
    rows     = block["rows"]
    if not headers:
        return
    col_count = max(len(headers), max((len(r) for r in rows), default=0))
    row_count = 1 + len(rows)

    tbl = slide.shapes.add_table(row_count, col_count, left, top, width, height).table

    # Header row
    for ci, h in enumerate(headers[:col_count]):
        cell = tbl.cell(0, ci)
        cell.text = h
        _style_cell(cell, bg=_DARK, fg=_WHITE, bold=True, size=Pt(10))

    # Data rows
    for ri, row_data in enumerate(rows, start=1):
        for ci in range(col_count):
            val  = row_data[ci] if ci < len(row_data) else ""
            cell = tbl.cell(ri, ci)
            cell.text = val
            bg = _LIGHT if ri % 2 == 0 else _WHITE
            _style_cell(cell, bg=bg, fg=_DARK, bold=False, size=Pt(9))


def _style_cell(cell, bg: str, fg: str, bold: bool, size) -> None:
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor.from_string(bg)
    for para in cell.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = RGBColor.from_string(fg)
            run.font.bold = bold
            run.font.size = size


# ── Drawing helpers ────────────────────────────────────────────────────────────

def _fill_background(slide, hex_color: str) -> None:
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(hex_color)


def _add_rect(slide, left, top, width, height, color: str) -> None:
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(color)
    shape.line.fill.background()


def _add_text_box(
    slide, text: str,
    left, top, width, height,
    font_size, bold: bool, color: str, align: str,
) -> None:
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.LEFT if align == "left" else PP_ALIGN.CENTER
    run = para.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)


# ── Path helper ────────────────────────────────────────────────────────────────

def _resolve_path(path: str, title: str, ext: str) -> Path:
    from xybernetex._models import _safe_filename  # local: only needed here
    p = Path(path)
    if p.is_dir():
        p = p / f"{_safe_filename(title)}{ext}"
    return p
